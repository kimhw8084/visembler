from __future__ import annotations

import base64
import binascii
import importlib.util
import io
import json
from pathlib import Path
from typing import Any, Mapping, Sequence

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt
from PIL import Image

from .files import validate_image_bytes, validate_pptx_bytes
from .domain import MODEL_MAX_BYTES, VisualizerContractError, canonical_model, stable_json

_VENDOR_ADAPTER = Path(__file__).with_name('vendor') / 'production_core' / 'tools' / 'ppt_template_adapter.py'
_MAX_ITEMS_PER_SLIDE = 12


def _adapter():
    spec=importlib.util.spec_from_file_location('company_ui_visualizer_frozen_ppt_adapter', _VENDOR_ADAPTER)
    if spec is None or spec.loader is None: raise RuntimeError('frozen Visualizer PPT adapter is unavailable')
    module=importlib.util.module_from_spec(spec); spec.loader.exec_module(module); return module


def _kind(entry: Mapping[str, Any]) -> str:
    engine=str(entry.get('engine') or '')
    if engine=='ImageMediaEngine': return 'image'
    if engine in {'CoreChartEngine','EngineeringChartEngine'}: return 'chart'
    if engine in {'TableEngine','MatrixEngine'}: return 'table'
    if engine in {'MetricEngine','ComparisonEngine'}: return 'kpi'
    if engine=='DiagramEngine': return 'diagram'
    if engine=='TimelineEngine': return 'timeline'
    if engine in {'TextEngine','EvidenceCompositeEngine','DecisionCompositeEngine','ProjectCompositeEngine'}: return 'text'
    return 'fallback'


def bound_export_items(model: Mapping[str, Any]) -> list[dict[str, Any]]:
    """Resolve canonical dataset bindings into the same export-facing fields as the editor.

    The report model remains untouched; this is an export projection so linked visuals
    always use their current shared source rather than stale item-local preview data.
    """
    datasets={str(dataset.get('id')):dataset for dataset in model.get('datasets') or [] if isinstance(dataset,Mapping)}
    resolved=[]
    for source in model.get('items') or []:
        if not isinstance(source,Mapping): continue
        entry=dict(source); dataset=datasets.get(str(entry.get('dataset_id') or ''))
        if not dataset:
            resolved.append(entry); continue
        fields=list(dataset.get('fields') or []); rows=[list(row) for row in dataset.get('rows') or [] if isinstance(row,Sequence) and not isinstance(row,(str,bytes))]
        mapping=entry.get('mapping') if isinstance(entry.get('mapping'),Mapping) else {}
        def index(role: str) -> int:
            field_id=mapping.get(role)
            return next((i for i,field in enumerate(fields) if isinstance(field,Mapping) and field.get('id')==field_id),-1)
        def first(role: str) -> Any:
            column=index(role)
            return next((row[column] for row in rows if column>=0 and column<len(row) and row[column] not in (None,'')),None)
        engine=str(entry.get('engine') or '')
        if engine=='TableEngine':
            entry['customTable']={'headers':[str(field.get('name') or field.get('id') or f'Column {i+1}') if isinstance(field,Mapping) else f'Column {i+1}' for i,field in enumerate(fields)],'rows':rows}
        elif engine=='MatrixEngine':
            entry['matrix']=[[str(field.get('name') or field.get('id') or f'Column {i+1}') if isinstance(field,Mapping) else f'Column {i+1}' for i,field in enumerate(fields)],*rows]
        elif engine=='TimelineEngine':
            label=index('category') if index('category')>=0 else index('x'); date=index('time')
            entry['milestones']=[{'label':str(row[label] if label>=0 and label<len(row) else i+1),'date':row[date] if date>=0 and date<len(row) else None} for i,row in enumerate(rows)]
        elif engine=='DiagramEngine':
            source_index,target_index=index('source'),index('target')
            edges=[(str(row[source_index]),str(row[target_index])) for row in rows if source_index>=0 and target_index>=0 and source_index<len(row) and target_index<len(row) and row[source_index] not in (None,'') and row[target_index] not in (None,'')]
            entry['edges']=edges; entry['nodes']=list(dict.fromkeys(node for edge in edges for node in edge))
        elif engine=='WaferFabEngine':
            x=index('die_x') if index('die_x')>=0 else index('x'); y=index('die_y') if index('die_y')>=0 else index('y'); value=index('value')
            entry['observations']=[{'x':row[x] if x>=0 and x<len(row) else None,'y':row[y] if y>=0 and y<len(row) else None,'value':row[value] if value>=0 and value<len(row) else None} for row in rows]
            entry.update({key:first(key) for key in ('wafer_id','lot_id','tool','chamber','recipe','process')})
        else:
            label=index('category') if index('category')>=0 else (index('time') if index('time')>=0 else index('x')); value=index('value') if index('value')>=0 else index('y')
            points=[(str(row[label] if label>=0 and label<len(row) else i+1),row[value] if value>=0 and value<len(row) else None) for i,row in enumerate(rows)]
            if engine in {'CoreChartEngine','EngineeringChartEngine'}: entry['data']=points;entry['rows']=[{'label':label,'value':value} for label,value in points];entry['observations']=[{'label':label,'value':value} for label,value in points]
            elif engine=='MetricEngine': entry['value']=next((value for _,value in reversed(points) if value is not None),None)
            elif engine=='ComparisonEngine':
                values=[value for _,value in points if value is not None];entry['before']=values[0] if values else None;entry['after']=values[-1] if values else None
        resolved.append(entry)
    return resolved


def _plan(model: Mapping[str, Any]) -> dict[str, Any]:
    report_items=list(model.get('items') or [])
    if not report_items: return _adapter().default_plan()
    positioned=[entry for entry in report_items if all(isinstance(entry.get(key),(int,float)) for key in ('x','y','w','h')) and entry['w']>0 and entry['h']>0]
    if len(positioned)==len(report_items):
        canvas_w=max(1200.0,max(float(entry['x'])+float(entry['w']) for entry in report_items))
        canvas_h=max(675.0,max(float(entry['y'])+float(entry['h']) for entry in report_items))
        return {'items':[{'kind':_kind(entry),'title':str(entry.get('title') or entry.get('element') or _kind(entry))[:100],
                          'nx':max(0.0,min(1.0,float(entry['x'])/canvas_w)),'ny':max(0.0,min(1.0,float(entry['y'])/canvas_h)),
                          'nw':max(.001,min(1.0,float(entry['w'])/canvas_w)),'nh':max(.001,min(1.0,float(entry['h'])/canvas_h))} for entry in report_items]}
    cols=1 if len(report_items)==1 else 2 if len(report_items)<=8 else 3
    rows=max(1,(len(report_items)+cols-1)//cols)
    gap_x=.025 if cols>1 else 0.0; gap_y=.035 if rows>1 else 0.0
    cell_w=(1-gap_x*(cols-1))/cols; cell_h=(1-gap_y*(rows-1))/rows
    items=[]
    for index,entry in enumerate(report_items):
        col=index%cols; row=index//cols
        items.append({'kind':_kind(entry),'title':str(entry.get('title') or entry.get('element') or _kind(entry))[:100],
                      'nx':col*(cell_w+gap_x),'ny':row*(cell_h+gap_y),'nw':cell_w,'nh':cell_h})
    return {'items':items}


def _export_pages(entries: list[dict[str, Any]]) -> list[list[dict[str, Any]]]:
    """Split deterministically, honoring optional integer page/page_index and page_break metadata."""
    explicit=any(isinstance(entry.get('page_index',entry.get('page')) ,int) for entry in entries)
    if explicit:
        grouped: dict[int,list[dict[str, Any]]]={}
        for entry in entries: grouped.setdefault(int(entry.get('page_index',entry.get('page',0))),[]).append(entry)
        return [grouped[index] for index in sorted(grouped)]
    ordered=sorted(entries,key=lambda entry:(int(entry.get('order',0)),str(entry.get('id',''))))
    pages: list[list[dict[str, Any]]]=[[]]
    for entry in ordered:
        if pages[-1] and (entry.get('page_break') is True or len(pages[-1])>=_MAX_ITEMS_PER_SLIDE): pages.append([])
        pages[-1].append(entry)
    return [page for page in pages if page]


def _display(value: Any) -> str:
    if value is None: return ''
    if isinstance(value,bool): return 'True' if value else 'False'
    return str(value)


def _semantic_text(entry: Mapping[str, Any]) -> str:
    engine=str(entry.get('engine') or '')
    if engine=='TextEngine': return str(entry.get('text') or entry.get('body') or '')
    if engine=='TimelineEngine':
        lines=[]
        for milestone in entry.get('milestones') or []:
            label=str(milestone.get('label') or '') if isinstance(milestone,Mapping) else str(milestone)
            date=milestone.get('date') if isinstance(milestone,Mapping) else None
            lines.append(f'{label} · {date}' if date not in (None,'') else label)
        return '\n'.join(lines)
    if engine=='DiagramEngine':
        nodes=[str(x) for x in entry.get('nodes') or []]; edges=entry.get('edges') or []
        return '\n'.join(nodes + [f'{edge[0]} → {edge[1]}' for edge in edges if isinstance(edge,Sequence) and not isinstance(edge,(str,bytes)) and len(edge)>=2])
    if engine in {'EvidenceCompositeEngine','DecisionCompositeEngine','ProjectCompositeEngine'}:
        return '\n'.join(x for x in [str(entry.get('statement') or ''),str(entry.get('detail') or ''),str(entry.get('status') or '')] if x)
    if engine=='ImageMediaEngine': return '\n'.join(x for x in [str(entry.get('caption') or ''),str(entry.get('alt') or '')] if x)
    if engine in {'WaferFabEngine','SmartLayoutEngine','InteractionLayer','EditorInfrastructure'}:
        return str(entry.get('configuration') or entry.get('behavior') or entry.get('route') or entry.get('element') or '')
    return str(entry.get('detail') or entry.get('element') or '')


def _chart_rows(entry: Mapping[str, Any]) -> list[tuple[str, Any]]:
    rows=entry.get('rows') or entry.get('data') or entry.get('observations') or []
    result=[]
    for i,row in enumerate(rows):
        if isinstance(row,Mapping): result.append((str(row.get('label') or row.get('x') or i+1),row.get('value')))
        elif isinstance(row,Sequence) and not isinstance(row,(str,bytes)) and len(row)>=2: result.append((str(row[0]),row[1]))
    return result or [('Value',None)]


def _table_grid(entry: Mapping[str, Any]) -> tuple[list[str],list[list[Any]]]:
    engine=str(entry.get('engine') or '')
    if engine=='MatrixEngine':
        rows=[list(r) for r in entry.get('matrix') or []]; cols=max((len(r) for r in rows),default=0)
        return [f'C{i+1}' for i in range(cols)],rows
    custom=entry.get('customTable') or {}
    headers=[str(x) for x in custom.get('headers') or []] if isinstance(custom,Mapping) else []
    rows=[list(r) for r in (custom.get('rows') or [])] if isinstance(custom,Mapping) else []
    if not rows: rows=[list(r) for r in entry.get('rows') or [] if isinstance(r,Sequence) and not isinstance(r,(str,bytes))]
    cols=max(len(headers),max((len(r) for r in rows),default=0),1)
    if not headers: headers=[f'Column {i+1}' for i in range(cols)]
    headers=(headers+['']*cols)[:cols]; rows=[(r+[None]*cols)[:cols] for r in rows]
    return headers,rows


def _set_semantic_metadata(shape: Any, entry: Mapping[str, Any]) -> None:
    """Embed the exact canonical element payload so 0, '0', blank and null remain distinguishable."""
    try:
        nodes=shape._element.xpath('.//p:cNvPr')
        if nodes: nodes[0].set('descr','VisualizerSemantic:'+json.dumps(entry,ensure_ascii=False,separators=(',',':')))
    except Exception:
        pass


def _set_report_metadata(slide: Any, model: Mapping[str, Any]) -> None:
    """Store bounded model context without changing the visible slide composition."""
    shape=slide.shapes.add_textbox(0,0,1,1)
    shape.name='VIZ::SemanticReport'
    nodes=shape._element.xpath('.//p:cNvPr')
    if nodes: nodes[0].set('descr','VisualizerSemanticReport:'+stable_json(model))


def _replace_image(slide: Any, shape: Any, entry: Mapping[str, Any], title: str, asset_data_url: Any=None) -> Any:
    src=entry.get('src') or (asset_data_url(entry['asset_id']) if asset_data_url and entry.get('asset_id') else None)
    if not isinstance(src,str) or not src.startswith('data:image/') or ';base64,' not in src:
        return shape
    _,encoded=src.split(',',1)
    raw=base64.b64decode(encoded,validate=True)
    image=Image.open(io.BytesIO(raw)); image.load()
    payload=io.BytesIO()
    image.convert('RGBA' if image.mode=='RGBA' else 'RGB').save(payload,format='PNG')
    left,top,width,height=shape.left,shape.top,shape.width,shape.height
    shape._element.getparent().remove(shape._element)
    picture=slide.shapes.add_picture(io.BytesIO(payload.getvalue()),left,top,width,height)
    picture.name=f'VIZ::{title} image'
    return picture


def _replace_diagram(slide: Any, shape: Any, entry: Mapping[str, Any], title: str) -> Any:
    nodes=[str(node) for node in entry.get('nodes') or []]
    if not nodes: return shape
    left,top,width,height=shape.left,shape.top,shape.width,shape.height
    shape._element.getparent().remove(shape._element)
    gap=Inches(.08); node_width=max(Inches(.55),(width-gap*(len(nodes)-1))//len(nodes)); node_height=max(Inches(.35),height//3)
    created={}
    for index,label in enumerate(nodes):
        node=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left+index*(node_width+gap),top+(height-node_height)//2,node_width,node_height)
        node.name=f'VIZ::{title}::{label}'; node.text_frame.text=label; created[label]=node
    for edge in entry.get('edges') or []:
        if not isinstance(edge,Sequence) or isinstance(edge,(str,bytes)) or len(edge)<2: continue
        source,target=created.get(str(edge[0])),created.get(str(edge[1]))
        if source is not None and target is not None: slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,source.left+source.width,source.top+source.height//2,target.left,target.top+target.height//2)
    return next(iter(created.values()))


def _fill_text_shape(shape: Any, entry: Mapping[str, Any], title: str) -> None:
    if not getattr(shape,'has_text_frame',False): return
    tf=shape.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=title; p.font.bold=True; p.font.size=Pt(13)
    body=_semantic_text(entry)
    if body:
        p2=tf.add_paragraph(); p2.text=body; p2.font.size=Pt(10)


def _fill_kpi(shape: Any, entry: Mapping[str, Any], title: str) -> None:
    if not getattr(shape,'has_text_frame',False): return
    value=entry.get('value')
    if value is None and str(entry.get('engine') or '')=='ComparisonEngine':
        value=f"{_display(entry.get('before'))} → {_display(entry.get('after'))}".strip()
    tf=shape.text_frame; tf.clear();p=tf.paragraphs[0];p.text=title;p.font.bold=True;p.font.size=Pt(10)
    p2=tf.add_paragraph();p2.text=_display(value);p2.font.bold=True;p2.font.size=Pt(24)
    unit=str(entry.get('unit') or '')
    if unit:
        p3=tf.add_paragraph();p3.text=unit;p3.font.size=Pt(9)


def _fill_chart(shape: Any, entry: Mapping[str, Any], title: str) -> None:
    if not getattr(shape,'has_chart',False): return
    rows=_chart_rows(entry);data=ChartData();data.categories=[r[0] for r in rows];data.add_series(title,[r[1] for r in rows]);shape.chart.replace_data(data)
    shape.chart.has_title=True;shape.chart.chart_title.text_frame.text=title;shape.chart.has_legend=False


def _replace_table(slide: Any, shape: Any, entry: Mapping[str, Any], title: str) -> Any:
    headers,rows=_table_grid(entry);cols=max(1,len(headers));nrows=max(2,1+len(rows))
    left,top,width,height=shape.left,shape.top,shape.width,shape.height
    shape._element.getparent().remove(shape._element)
    new_shape=slide.shapes.add_table(nrows,cols,left,top,width,height);new_shape.name=f'VIZ::{title}'
    table=new_shape.table
    for c,h in enumerate(headers): table.cell(0,c).text=_display(h)
    for r,row in enumerate(rows,1):
        for c in range(cols): table.cell(r,c).text=_display(row[c] if c<len(row) else None)
    for r in range(nrows):
        for c in range(cols):
            for p in table.cell(r,c).text_frame.paragraphs:p.font.size=Pt(9);p.font.bold=(r==0)
    return new_shape


def _apply_semantics(slide: Any, before_count: int, entries: list[Mapping[str, Any]], plan: Mapping[str, Any], asset_data_url: Any=None) -> None:
    created=[slide.shapes[i] for i in range(before_count,len(slide.shapes))]
    # adapter emits one top-level shape per plan item; replacement tables are processed from the correlated originals.
    if len(created)<len(entries): raise RuntimeError('frozen PowerPoint adapter created fewer shapes than planned')
    for index,(entry,item) in enumerate(zip(entries,plan['items'])):
        shape=created[index];kind=item['kind'];title=item['title']
        if kind=='image': shape=_replace_image(slide,shape,entry,title,asset_data_url)
        elif kind=='diagram': shape=_replace_diagram(slide,shape,entry,title)
        elif kind=='timeline': _fill_text_shape(shape,entry,title)
        elif kind=='kpi': _fill_kpi(shape,entry,title)
        elif kind=='chart': _fill_chart(shape,entry,title)
        elif kind=='table': shape=_replace_table(slide,shape,entry,title)
        elif kind=='fallback': _fill_text_shape(shape,{**entry,'detail':f'Controlled fallback · {entry.get("element") or "specialized visual"} remains semantic metadata; recreate this visual natively in Visembler.'},title)
        else: _fill_text_shape(shape,entry,title)
        _set_semantic_metadata(shape,entry)


def export_pptx(template_bytes: bytes | None, model: Mapping[str, Any], *, slide_index: int = 0, placeholder: str = 'VISUALIZER_CONTENT', asset_data_url: Any=None) -> bytes:
    """Export the authored report into an optional template or a clean blank deck."""
    if template_bytes is None:
        prs=Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
    else:
        validate_pptx_bytes(template_bytes)
        prs=Presentation(io.BytesIO(template_bytes))
    if slide_index < 0 or slide_index >= len(prs.slides): raise ValueError('PPT slide index is out of range')
    try: semantic_model=canonical_model(model)
    except VisualizerContractError: semantic_model=None
    if semantic_model is not None and asset_data_url:
        semantic_model=json.loads(stable_json(semantic_model))
        for entry in semantic_model['items']:
            if entry.get('engine')=='ImageMediaEngine' and entry.get('asset_id'):
                entry['src']=asset_data_url(entry['asset_id']); entry.pop('asset_id',None)
    entries=bound_export_items(semantic_model or model); adapter=_adapter()
    for page_index,page_entries in enumerate(_export_pages(entries)):
        slide=prs.slides[slide_index] if page_index==0 else prs.slides.add_slide(prs.slide_layouts[6]); target_slide_index=slide_index if page_index==0 else len(prs.slides)-1
        plan=_plan({**model,'items':page_entries}); before_count=len(slide.shapes)
        adapter.insert(prs,slide_index=target_slide_index,placeholder=placeholder if page_index==0 else None,plan=plan)
        _apply_semantics(slide,before_count,page_entries,plan,asset_data_url)
        if page_index==0 and semantic_model is not None: _set_report_metadata(slide,semantic_model)
    output=io.BytesIO(); prs.save(output); payload=output.getvalue(); validate_pptx_bytes(payload); return payload


def import_visembler_pptx(payload: bytes) -> dict[str, Any] | None:
    """Rebuild a report only from exact exported metadata; never infer from ordinary shapes."""
    validate_pptx_bytes(payload)
    prs=Presentation(io.BytesIO(payload)); items=[]; seen=set(); report_context=None
    for slide in prs.slides:
        for shape in slide.shapes:
            descriptions=[]
            try: descriptions=[node.get('descr') for node in shape._element.xpath('.//p:cNvPr')]
            except Exception: descriptions=[]
            for description in descriptions:
                if isinstance(description,str) and description.startswith('VisualizerSemanticReport:'):
                    if report_context is not None: raise VisualizerContractError('semantic PowerPoint contains multiple report payloads')
                    encoded=description.removeprefix('VisualizerSemanticReport:')
                    if len(encoded.encode('utf-8'))>MODEL_MAX_BYTES: raise VisualizerContractError('semantic PowerPoint payload exceeds model limit')
                    try: report_context=json.loads(encoded,parse_constant=lambda _value: (_ for _ in ()).throw(ValueError('non-finite value')))
                    except (json.JSONDecodeError,ValueError) as exc: raise VisualizerContractError('malformed Visembler report payload') from exc
                    if not isinstance(report_context,Mapping): raise VisualizerContractError('semantic Visembler report payload must be an object')
                    continue
                if not isinstance(description,str) or not description.startswith('VisualizerSemantic:'): continue
                encoded=description.removeprefix('VisualizerSemantic:')
                if len(encoded.encode('utf-8'))>MODEL_MAX_BYTES: raise VisualizerContractError('semantic PowerPoint payload exceeds model limit')
                try: entry=json.loads(encoded,parse_constant=lambda _value: (_ for _ in ()).throw(ValueError('non-finite value')))
                except (json.JSONDecodeError,ValueError) as exc: raise VisualizerContractError('malformed VisualizerSemantic payload') from exc
                if not isinstance(entry,Mapping): raise VisualizerContractError('semantic PowerPoint payload must be an object')
                entry=dict(entry); entry_id=str(entry.get('id') or '')
                if not entry_id or entry_id in seen: raise VisualizerContractError('semantic PowerPoint contains duplicate or missing element ids')
                # Older geometry-less exports can still recover a deterministic canvas position.
                if not all(isinstance(entry.get(key),(int,float)) for key in ('x','y','w','h')):
                    entry.update({'x':shape.left/prs.slide_width*1200,'y':shape.top/prs.slide_height*675,'w':shape.width/prs.slide_width*1200,'h':shape.height/prs.slide_height*675})
                seen.add(entry_id);items.append(entry)
    if not items:
        if report_context is not None: raise VisualizerContractError('semantic Visembler report is missing element payloads')
        return None
    for entry in items:
        if entry.get('engine')!='ImageMediaEngine' or not entry.get('src'): continue
        src=entry.get('src')
        if not isinstance(src,str) or ';base64,' not in src: raise VisualizerContractError('semantic PowerPoint image must be an embedded PNG, JPEG, or WebP')
        prefix,encoded=src.split(',',1)
        if prefix not in {'data:image/png;base64','data:image/jpeg;base64','data:image/webp;base64'}: raise VisualizerContractError('semantic PowerPoint image format is unsupported')
        try: validate_image_bytes(base64.b64decode(encoded,validate=True))
        except (ValueError,binascii.Error,VisualizerContractError) as exc: raise VisualizerContractError('semantic PowerPoint image is invalid') from exc
    # Element metadata preserves the bound export projection, not its source dataset.
    # Drop dangling links so imported elements stay canonical and immediately editable.
    for entry in items:
        if entry.pop('dataset_id',None) is not None: entry.pop('mapping',None)
    next_ids=[int(str(item['id'])[1:])+1 for item in items if str(item['id']).startswith('c') and str(item['id'])[1:].isdigit()]
    model=canonical_model({'items':items,'nextId':max([1,*next_ids])}) if report_context is None else canonical_model(report_context)
    if report_context is not None and {str(item.get('id')) for item in model['items']} != seen:
        raise VisualizerContractError('semantic Visembler report element payloads do not match report context')
    if len(stable_json(model).encode('utf-8'))>MODEL_MAX_BYTES: raise VisualizerContractError('semantic PowerPoint report exceeds model limit')
    return model
