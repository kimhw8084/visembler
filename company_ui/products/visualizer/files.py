from __future__ import annotations

import io, os, re, zipfile
from pathlib import PurePosixPath
from typing import Iterable

from PIL import Image, UnidentifiedImageError

from .domain import IMAGE_EMBED_MAX_BYTES, VisualizerContractError

PPT_MAX_BYTES=100*1024*1024
PPT_MAX_ENTRIES=5000
PPT_MAX_EXPANDED=350*1024*1024
PPT_MAX_RATIO=250.0
IMAGE_MAX_PIXELS=20_000_000
_ALLOWED_IMAGE_FORMATS={'PNG','JPEG','WEBP'}


def validate_image_bytes(data: bytes) -> dict[str, object]:
    if not data: raise VisualizerContractError('image is empty')
    if len(data)>IMAGE_EMBED_MAX_BYTES: raise VisualizerContractError(f'image exceeds {IMAGE_EMBED_MAX_BYTES} bytes')
    try:
        with Image.open(io.BytesIO(data)) as im:
            fmt=(im.format or '').upper(); width,height=im.size
            if fmt not in _ALLOWED_IMAGE_FORMATS: raise VisualizerContractError(f'unsupported image format: {fmt or "unknown"}')
            if width<1 or height<1 or width*height>IMAGE_MAX_PIXELS: raise VisualizerContractError('image dimensions exceed safety limit')
            im.verify()
    except (UnidentifiedImageError,OSError) as exc:
        raise VisualizerContractError('invalid image content') from exc
    mime={'PNG':'image/png','JPEG':'image/jpeg','WEBP':'image/webp'}[fmt]
    return {'format':fmt,'mime':mime,'width':width,'height':height,'bytes':len(data)}


def _safe_zip_name(name: str) -> bool:
    if not name or '\\' in name or '\x00' in name: return False
    path=PurePosixPath(name)
    if path.is_absolute() or any(part in {'','..'} for part in path.parts): return False
    return True


def validate_pptx_bytes(data: bytes) -> dict[str, int]:
    if not data: raise VisualizerContractError('PowerPoint is empty')
    if len(data)>PPT_MAX_BYTES: raise VisualizerContractError(f'PowerPoint exceeds {PPT_MAX_BYTES} bytes')
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            infos=zf.infolist()
            if len(infos)>PPT_MAX_ENTRIES: raise VisualizerContractError('PowerPoint has too many archive entries')
            seen=set(); expanded=0
            for info in infos:
                name=info.filename
                key=name.casefold()
                if not _safe_zip_name(name): raise VisualizerContractError(f'unsafe PowerPoint archive path: {name!r}')
                if key in seen: raise VisualizerContractError(f'duplicate PowerPoint archive path: {name!r}')
                seen.add(key); expanded += int(info.file_size)
                if expanded>PPT_MAX_EXPANDED: raise VisualizerContractError('PowerPoint expanded size exceeds safety limit')
                if info.file_size and info.compress_size:
                    ratio=info.file_size/max(1,info.compress_size)
                    if ratio>PPT_MAX_RATIO and info.file_size>1_000_000: raise VisualizerContractError('PowerPoint archive contains suspicious compression ratio')
            required={'[content_types].xml','ppt/presentation.xml'}
            names={i.filename.casefold() for i in infos}
            if not required.issubset(names): raise VisualizerContractError('file is not a valid PPTX package')
    except zipfile.BadZipFile as exc:
        raise VisualizerContractError('PowerPoint is not a valid ZIP package') from exc
    from pptx import Presentation
    try: prs=Presentation(io.BytesIO(data))
    except Exception as exc: raise VisualizerContractError('PowerPoint package cannot be opened') from exc
    if not prs.slides: raise VisualizerContractError('PowerPoint contains no slides')
    return {'bytes':len(data),'entries':len(infos),'expanded_bytes':expanded,'slides':len(prs.slides)}
