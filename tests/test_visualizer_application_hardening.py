from __future__ import annotations

import io
import json
import os
import zipfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from company_ui.products.visualizer.domain import RevisionConflictError, VisualizerContractError, canonical_model
from company_ui.products.visualizer.files import validate_image_bytes, validate_pptx_bytes
from company_ui.products.visualizer.ppt_service import export_pptx, import_visembler_pptx
from company_ui.products.visualizer.page import MAX_PRESETS, _decode_bridge_event, _normalize_presets, _report_options, _validate_model_images
from company_ui.products.visualizer.repository import ReportRepository
from company_ui.products.visualizer.runtime import application_environment, resolve_runtime
from company_ui.products.visualizer.templates import REPORT_TEMPLATES, template_model


def test_runtime_secret_is_generated_reused_and_reaches_nicegui_kwargs(tmp_path: Path):
    base={'COMPANY_UI_ENVIRONMENT':'dev','COMPANY_UI_VISUALIZER_DATA_DIR':str(tmp_path),'HOST':'127.0.0.1','PORT':'8123'}
    first=application_environment(base); second=application_environment(base); config, env=resolve_runtime(base); kwargs=config.nicegui_run_kwargs(env)
    assert first['COMPANY_UI_STORAGE_SECRET']==second['COMPANY_UI_STORAGE_SECRET']
    assert len(first['COMPANY_UI_STORAGE_SECRET'])>=32
    assert kwargs['storage_secret']==first['COMPANY_UI_STORAGE_SECRET']
    assert kwargs['host']=='127.0.0.1' and kwargs['port']==8123
    assert (tmp_path/'.storage_secret').read_text().strip()==first['COMPANY_UI_STORAGE_SECRET']


def test_concurrent_local_startups_share_one_storage_secret(tmp_path: Path):
    from concurrent.futures import ThreadPoolExecutor
    base={'COMPANY_UI_ENVIRONMENT':'dev','COMPANY_UI_VISUALIZER_DATA_DIR':str(tmp_path)}
    with ThreadPoolExecutor(max_workers=8) as pool:
        values=list(pool.map(lambda _: application_environment(base)['COMPANY_UI_STORAGE_SECRET'], range(24)))
    assert len(set(values))==1
    assert (tmp_path/'.storage_secret').read_text().strip()==values[0]


def test_production_missing_storage_secret_fails_closed(tmp_path: Path):
    with pytest.raises(RuntimeError, match='COMPANY_UI_STORAGE_SECRET is required'):
        application_environment({'COMPANY_UI_ENVIRONMENT':'prod','COMPANY_UI_VISUALIZER_DATA_DIR':str(tmp_path)})


def test_repository_revision_idempotency_cleanup_and_corruption_isolation(tmp_path: Path):
    repo=ReportRepository(tmp_path)
    r=repo.create('r1',model=template_model('blank'))
    model=canonical_model({'items':[{'id':'c1','type':'text','order':0,'text':'a'}],'nextId':2})
    committed=repo.commit('r1',base_revision=1,model=model,commit_id='same')
    assert repo.commit('r1',base_revision=1,model=model,commit_id='same').revision==committed.revision
    with pytest.raises(RevisionConflictError): repo.commit('r1',base_revision=1,model=model,commit_id='new')
    assert repo.delete_if_blank('r1',expected_revision=committed.revision) is False
    repo.create('bad',model=template_model('blank')); (tmp_path/'bad.json').write_text('{bad',encoding='utf-8')
    assert [x.report_id for x in repo.list()]==['r1']
    assert any((tmp_path/'_quarantine').glob('bad.*.corrupt.json'))


def test_two_repository_instances_cannot_accept_same_base_revision(tmp_path: Path):
    from concurrent.futures import ThreadPoolExecutor
    repo=ReportRepository(tmp_path); repo.create('shared',model=template_model('blank'))
    model_a=canonical_model({'items':[{'id':'a','type':'text','order':0,'text':'A'}],'nextId':2})
    model_b=canonical_model({'items':[{'id':'b','type':'text','order':0,'text':'B'}],'nextId':2})
    def write(which):
        local=ReportRepository(tmp_path)
        try:
            value=model_a if which=='a' else model_b
            return ('ok',local.commit('shared',base_revision=1,model=value,commit_id=which).revision)
        except RevisionConflictError as exc:
            return ('conflict',exc.expected)
    with ThreadPoolExecutor(max_workers=2) as pool:
        results=list(pool.map(write,['a','b']))
    assert sorted(kind for kind,_ in results)==['conflict','ok']
    final=repo.get('shared'); assert final.revision==2 and len(final.model['items'])==1


def test_repository_history_checkpoint_restore_duplicate_retention_and_trash_recovery(tmp_path: Path):
    repo=ReportRepository(tmp_path,history_limit=4)
    first=repo.create('history',title='History',model=canonical_model({'items':[{'id':'c1','type':'text','order':0,'text':'one'}]}))
    second=repo.commit('history',base_revision=first.revision,model=canonical_model({'items':[{'id':'c1','type':'text','order':0,'text':'two'}]}),commit_id='two')
    checkpoint=repo.checkpoint('history','Before restore',expected_revision=second.revision)
    restored=repo.restore_history('history','r1',expected_revision=second.revision)
    assert restored.revision==3 and restored.model['items'][0]['text']=='one'
    copied=repo.duplicate_from_history('history',checkpoint['history_id'],'history-copy',title='Historical copy')
    assert copied.model['items'][0]['text']=='two'
    assert len(repo.list_history('history'))<=4
    repo.trash_report('history',expected_revision=restored.revision)
    assert not repo._history_path('history').exists() and repo._trash_history_path('history').exists()
    repo.restore('history'); assert repo._history_path('history').exists()


def test_history_corruption_isolated_and_concurrent_commits_keep_a_single_revision(tmp_path: Path):
    from concurrent.futures import ThreadPoolExecutor
    repo=ReportRepository(tmp_path); first=repo.create('history',model=canonical_model({'items':[{'id':'c1','type':'text','order':0,'text':'one'}]}))
    def commit(text: str):
        local=ReportRepository(tmp_path)
        try: return local.commit('history',base_revision=first.revision,model=canonical_model({'items':[{'id':'c1','type':'text','order':0,'text':text}]}),commit_id=text).revision
        except RevisionConflictError: return None
    with ThreadPoolExecutor(max_workers=2) as pool: outcomes=list(pool.map(commit,['two','three']))
    assert sorted(value is None for value in outcomes)==[False,True]
    broken=repo._history_path('history')/'broken.json'; broken.write_text('{bad',encoding='utf-8')
    assert repo.list_history('history') and any(repo.quarantine.glob('*.history-corrupt.json'))


def test_conflicted_local_transaction_can_be_safely_reapplied_on_latest_revision(tmp_path: Path):
    repo=ReportRepository(tmp_path); repo.create('shared',model=template_model('blank'))
    remote=canonical_model({'items':[{'id':'remote','type':'text','order':0,'text':'remote'}],'nextId':3})
    local=canonical_model({'items':[{'id':'local','type':'text','order':0,'text':'local'}],'nextId':3})
    repo.commit('shared',base_revision=1,model=remote,commit_id='remote-1')
    with pytest.raises(RevisionConflictError): repo.commit('shared',base_revision=1,model=local,commit_id='local-stale')
    latest=repo.get('shared')
    recovered=canonical_model({**latest.model,'items':[*latest.model['items'],local['items'][0]],'nextId':3})
    committed=repo.commit('shared',base_revision=latest.revision,model=recovered,commit_id='local-reapply-1')
    assert committed.revision==3
    assert [entry['id'] for entry in committed.model['items']]==['remote','local']
    assert repo.commit('shared',base_revision=latest.revision,model=recovered,commit_id='local-reapply-1').revision==3


def test_report_trash_restore_and_duplicate_title_labels_are_recoverable(tmp_path: Path):
    repo=ReportRepository(tmp_path)
    first=repo.create('first',title='Untitled report',model=template_model('blank'))
    second=repo.create('second',title='Untitled report',model=template_model('blank'))
    labels=_report_options(repo)
    assert len(labels)==2 and all('New blank report' in label for label in labels.values())
    trashed=repo.trash_report(first.report_id,expected_revision=first.revision)
    assert trashed.report_id=='first' and [record.report_id for record in repo.list()]==['second']
    assert [record.report_id for record in repo.list_trash()]==['first']
    restored=repo.restore('first')
    assert restored.report_id=='first' and len(repo.list())==2 and repo.list_trash()==[]


def test_templates_are_deep_copied_and_blank_is_really_blank():
    blank=template_model('blank'); assert blank['items']==[] and blank['groups']=={} and blank['mode']=='guided'
    for key in REPORT_TEMPLATES:
        one=template_model(key); two=template_model(key); one['items'][0]['title']='changed'; assert two['items'][0]['title']!='changed'


def test_preset_storage_limit_duplicate_names_and_reload_normalization_are_deterministic():
    model=template_model('blank')
    raw=[{'id':f'p{index:02}', 'name':'Review' if index<2 else f'Preset {index:02}', 'model':model} for index in range(MAX_PRESETS+5)]
    normalized=_normalize_presets(raw)
    assert len(normalized)==MAX_PRESETS
    assert [entry['name'] for entry in normalized[:2]]==['Review','Review']
    assert _normalize_presets(json.loads(json.dumps(normalized)))==normalized


def _image_bytes(fmt: str='PNG') -> bytes:
    out=io.BytesIO(); Image.new('RGB',(20,10),(10,20,30)).save(out,format=fmt); return out.getvalue()


@pytest.mark.parametrize('fmt', ['PNG','JPEG','WEBP'])
def test_image_content_validation_allowlist(fmt: str):
    result=validate_image_bytes(_image_bytes(fmt)); assert result['format']==fmt and result['width']==20 and result['height']==10


def test_image_rejects_active_or_invalid_content():
    with pytest.raises(VisualizerContractError): validate_image_bytes(b'<svg onload="alert(1)"></svg>')
    with pytest.raises(VisualizerContractError): validate_image_bytes(b'not-an-image')


def test_embedded_image_server_validation_rejects_external_url_and_accepts_data_url():
    raw=_image_bytes(); import base64
    ok=canonical_model({'items':[{'id':'c1','type':'image','engine':'ImageMediaEngine','order':0,'src':'data:image/png;base64,'+base64.b64encode(raw).decode()}]})
    _validate_model_images(ok)
    bad=canonical_model({'items':[{'id':'c1','type':'image','engine':'ImageMediaEngine','order':0,'src':'https://example.invalid/a.png'}]})
    with pytest.raises(VisualizerContractError): _validate_model_images(bad)


def test_repository_migrates_deduplicates_and_garbage_collects_image_assets(tmp_path: Path):
    import base64
    repo=ReportRepository(tmp_path); image='data:image/png;base64,'+base64.b64encode(_image_bytes()).decode()
    model=lambda item_id: canonical_model({'items':[{'id':item_id,'type':'image','engine':'ImageMediaEngine','order':0,'src':image}]})
    first=repo.create('first',model=model('c1')); second=repo.create('second',model=model('c2'))
    asset_id=first.model['items'][0]['asset_id']
    assert 'src' not in first.model['items'][0] and second.model['items'][0]['asset_id']==asset_id
    assert repo.assets.ids()=={asset_id}
    repo.delete('first',expected_revision=first.revision); assert repo.assets.ids()=={asset_id}
    repo.delete('second',expected_revision=second.revision); assert repo.assets.ids()==set()


def test_repository_asset_corruption_and_ppt_export_recovery_are_explicit(tmp_path: Path):
    import base64
    repo=ReportRepository(tmp_path); image='data:image/png;base64,'+base64.b64encode(_image_bytes()).decode()
    record=repo.create('image',model=canonical_model({'items':[{'id':'c1','type':'image','engine':'ImageMediaEngine','order':0,'src':image}]}))
    output=export_pptx(_pptx(),record.model,asset_data_url=repo.assets.data_url)
    assert any(getattr(shape,'shape_type',None)==13 for shape in Presentation(io.BytesIO(output)).slides[0].shapes)
    repo.assets._path(record.model['items'][0]['asset_id']).write_bytes(b'corrupt')
    with pytest.raises(VisualizerContractError,match='corrupt'):
        repo.get('image')


def _pptx() -> bytes:
    prs=Presentation(); prs.slides.add_slide(prs.slide_layouts[6]); out=io.BytesIO(); prs.save(out); return out.getvalue()


def test_pptx_validation_and_archive_traversal_rejection():
    data=_pptx(); assert validate_pptx_bytes(data)['slides']==1
    out=io.BytesIO()
    with zipfile.ZipFile(out,'w',zipfile.ZIP_DEFLATED) as z:
        z.writestr('../evil','x'); z.writestr('[Content_Types].xml','x'); z.writestr('ppt/presentation.xml','x')
    with pytest.raises(VisualizerContractError, match='unsafe PowerPoint archive path'): validate_pptx_bytes(out.getvalue())


def test_bridge_rejects_unknown_version_type_and_oversize():
    good={'bridge_version':1,'type':'report.save_requested','payload':{}}
    assert _decode_bridge_event(json.dumps(good))['type']=='report.save_requested'
    with pytest.raises(VisualizerContractError): _decode_bridge_event({'bridge_version':2,'type':'report.save_requested','payload':{}})
    with pytest.raises(VisualizerContractError): _decode_bridge_event({'bridge_version':1,'type':'pointermove','payload':{}})


def test_preset_normalization_isolates_corruption_and_preserves_zero():
    raw=[{'id':'good','name':'Good','model':{'items':[{'id':'c1','type':'metric','order':0,'value':0}]}},{'name':'bad','model':{'items':[{'id':'x','order':0}]}}]
    result=_normalize_presets(raw); assert len(result)==1 and result[0]['model']['items'][0]['value']==0


def test_image_pixel_bomb_is_rejected_even_when_compressed_small():
    out=io.BytesIO(); Image.new('RGB',(5000,5000),(0,0,0)).save(out,format='PNG',compress_level=9)
    assert len(out.getvalue()) < 750_000
    with pytest.raises(VisualizerContractError,match='dimensions exceed safety limit'):
        validate_image_bytes(out.getvalue())


def test_pptx_duplicate_path_and_compression_bomb_rejected():
    duplicate=io.BytesIO()
    with zipfile.ZipFile(duplicate,'w',zipfile.ZIP_DEFLATED) as z:
        z.writestr('[Content_Types].xml','x'); z.writestr('ppt/presentation.xml','x'); z.writestr('ppt/presentation.xml','again')
    with pytest.raises(VisualizerContractError,match='duplicate PowerPoint archive path'):
        validate_pptx_bytes(duplicate.getvalue())
    bomb=io.BytesIO()
    with zipfile.ZipFile(bomb,'w',zipfile.ZIP_DEFLATED,compresslevel=9) as z:
        z.writestr('[Content_Types].xml','x'); z.writestr('ppt/presentation.xml','x'); z.writestr('ppt/media/zeros.bin',b'0'*(2*1024*1024))
    with pytest.raises(VisualizerContractError,match='suspicious compression ratio'):
        validate_pptx_bytes(bomb.getvalue())


def test_editable_ppt_round_trip_remains_openable_and_adds_native_shapes():
    import base64
    template=_pptx(); before=Presentation(io.BytesIO(template)); before_count=len(before.slides[0].shapes)
    image='data:image/png;base64,'+base64.b64encode(_image_bytes()).decode()
    model=canonical_model({'items':[
        {'id':'m1','type':'metric','engine':'MetricEngine','order':0,'title':'Yield','value':0,'x':0,'y':0,'w':300,'h':120},
        {'id':'t1','type':'text','engine':'TextEngine','order':1,'title':'Finding','text':'Native editable text','x':600,'y':300,'w':400,'h':120},
        {'id':'i1','type':'image','engine':'ImageMediaEngine','order':2,'title':'Wafer image','src':image,'x':100,'y':400,'w':240,'h':180},
    ],'nextId':4})
    output=export_pptx(template,model); validate_pptx_bytes(output)
    after=Presentation(io.BytesIO(output)); assert len(after.slides)==1 and len(after.slides[0].shapes)>before_count
    assert any(getattr(shape,'has_text_frame',False) for shape in after.slides[0].shapes)
    picture=next(shape for shape in after.slides[0].shapes if getattr(shape,'shape_type',None)==13)
    assert picture.width>0 and picture.height>0
    semantic=[shape for shape in after.slides[0].shapes if getattr(shape,'name','').startswith('VIZ::Yield')][0]
    assert semantic.left < after.slide_width/2


def test_ppt_export_without_an_uploaded_template_creates_an_openable_blank_deck():
    model=canonical_model({'items':[{'id':'t1','type':'text','engine':'TextEngine','order':0,'title':'Finding','text':'Export without setup','x':0,'y':0,'w':600,'h':180}],'nextId':2})
    output=export_pptx(None,model)
    validate_pptx_bytes(output)
    assert len(Presentation(io.BytesIO(output)).slides)==1


def test_ppt_export_paginates_deterministically_and_keeps_mixed_semantics():
    template=_pptx()
    items=[]
    for index in range(14):
        engine=('TimelineEngine' if index==0 else 'DiagramEngine' if index==1 else 'MetricEngine')
        item={'id':f'c{index}','type':'text','engine':engine,'element':'Event Timeline' if engine=='TimelineEngine' else 'Data Flow' if engine=='DiagramEngine' else 'Hero KPI','order':index,'title':f'Item {index}','value':0}
        if engine=='TimelineEngine': item['milestones']=[{'label':'Collect','date':'2026-01-01'},{'label':'Verify','date':'2026-01-02'}]
        if engine=='DiagramEngine': item.update(nodes=['Etch','Inspect'],edges=[['Etch','Inspect']])
        items.append(item)
    output=export_pptx(template,canonical_model({'items':items,'nextId':20}))
    deck=Presentation(io.BytesIO(output)); assert len(deck.slides)==2
    metadata=[]
    for slide in deck.slides:
        for shape in slide.shapes:
            for node in shape._element.xpath('.//p:cNvPr'):
                if str(node.get('descr','')).startswith('VisualizerSemantic:'): metadata.append(node.get('descr'))
    assert len(metadata)==14 and any('Event Timeline' in value for value in metadata) and any('"value":0' in value for value in metadata)


def test_visembler_pptx_round_trip_recovers_exact_semantic_values_and_geometry():
    import base64
    template=_pptx()
    image='data:image/png;base64,'+base64.b64encode(_image_bytes()).decode()
    source=canonical_model({'items':[
        {'id':'c1','type':'metric','engine':'MetricEngine','order':0,'title':'Zero','value':0,'x':12,'y':24,'w':220,'h':100},
        {'id':'c2','type':'text','engine':'TextEngine','order':1,'title':'String zero','text':'0','x':300,'y':24,'w':220,'h':100},
        {'id':'c3','type':'text','engine':'TextEngine','order':2,'title':'Blank','text':'','detail':None,'x':12,'y':180,'w':220,'h':100},
        {'id':'c4','type':'image','engine':'ImageMediaEngine','order':3,'title':'Image','src':image,'x':300,'y':180,'w':220,'h':100},
    ],'nextId':5})
    restored=import_visembler_pptx(export_pptx(template,source))
    assert restored is not None
    by_id={item['id']:item for item in restored['items']}
    assert by_id['c1']['value']==0 and by_id['c2']['text']=='0' and by_id['c3']['text']=='' and by_id['c3']['detail'] is None
    assert by_id['c1']['x']==12 and by_id['c1']['h']==100
    assert by_id['c4']['src']==image
    assert restored==source
    assert import_visembler_pptx(template) is None


def test_visembler_pptx_import_rejects_malformed_metadata_and_ignores_normal_shapes():
    prs=Presentation(); slide=prs.slides.add_slide(prs.slide_layouts[6]); shape=slide.shapes.add_textbox(0,0,100,100)
    output=io.BytesIO(); prs.save(output)
    assert import_visembler_pptx(output.getvalue()) is None
    shape._element.xpath('.//p:cNvPr')[0].set('descr','VisualizerSemantic:{not-json')
    output=io.BytesIO(); prs.save(output)
    with pytest.raises(VisualizerContractError,match='malformed VisualizerSemantic payload'):
        import_visembler_pptx(output.getvalue())


def test_visembler_pptx_import_preserves_bound_data_and_dataset_linkage():
    template=_pptx()
    source=canonical_model({'datasets':[{'id':'d1','fields':[{'id':'lot','name':'Lot'},{'id':'yield','name':'Yield'}],'rows':[['L01',0],['L02',None]]}], 'items':[
        {'id':'c1','type':'table','engine':'TableEngine','order':0,'dataset_id':'d1','mapping':{'category':'lot','value':'yield'},'x':0,'y':0,'w':300,'h':200},
    ],'nextId':2})
    restored=import_visembler_pptx(export_pptx(template,source))
    assert restored==source
    item=restored['items'][0]
    assert item['dataset_id']=='d1' and item['mapping']=={'category':'lot','value':'yield'}


def test_visualizer_page_does_not_touch_tab_storage_before_client_connection():
    from company_ui.products.visualizer import page as visualizer_page_module
    source=Path(visualizer_page_module.__file__).read_text(encoding='utf-8')
    assert 'NiceGUIStateServices.tab_store()' not in source
    assert 'NiceGUIStateServices.user_store()' in source
