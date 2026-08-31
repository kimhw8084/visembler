from __future__ import annotations

import hashlib
import importlib
import io
import json
import tomllib
from importlib.metadata import version
from pathlib import Path

import pytest
from pptx import Presentation

from company_ui.products.visualizer.ppt_service import export_pptx

ROOT=Path(__file__).resolve().parents[1]
PRODUCT=ROOT/'company_ui/products/visualizer'
GOLDEN='d8ebd4378f01b7c52a7a4be57c578c22adf29b899cc08a370cf084881195343e'


def _pptx() -> bytes:
    prs=Presentation(); prs.slides.add_slide(prs.slide_layouts[6]); out=io.BytesIO(); prs.save(out); return out.getvalue()


def test_r14_capability_matrix_has_product_level_contract_for_every_element():
    data=json.loads((PRODUCT/'contracts/ELEMENT_CAPABILITY_MATRIX.json').read_text())
    assert data['schema_version']==2 and data['capability_release']=='R14'
    assert data['count']==248 and data['engines']==17 and len(data['rows'])==248
    required={'element_id','semantic_family','visual_grammar','default_model_ref','library_thumbnail','inspector_schema','direct_edit_behavior','paste_behavior','empty_state','smart_layout_constraints','geometry_behavior','save_reload','undo_redo','responsive_behavior','accessibility_behavior','export_eligibility'}
    for row in data['rows']:
        assert required <= row.keys()
        smart=row['smart_layout_constraints']
        assert smart['min_width']>0 and smart['min_height']>0
        assert smart['overflow_policy']=='reflow_or_actionable_conflict'
        assert row['geometry_behavior']=={'smart':'semantic_owned','guided':'user_owned_snapped','free':'user_owned_unsnapped','minimum_enforced':True}
    assert len({row['inspector_schema'] for row in data['rows']})>=25
    assert len({row['library_thumbnail'] for row in data['rows']})>=150


def test_r14_metric_variants_have_distinct_semantic_inspector_schemas():
    rows={r['element']:r for r in json.loads((PRODUCT/'contracts/ELEMENT_CAPABILITY_MATRIX.json').read_text())['rows']}
    names=['Metric Ring','Metric Ladder','Confidence Metric','Capacity Metric','Rate Metric','Threshold Metric','Metric with Sparkline']
    assert len({rows[name]['inspector_schema'] for name in names})==len(names)


def test_r14_editor_exposes_semantic_mode_and_transient_guide_contracts():
    js=(PRODUCT/'assets/integrated_editor.mjs').read_text()
    assert "dragState.active == false" not in js  # no comment-only certification shortcut
    assert 'clearTransientInteractionVisuals(`pointer-${kind}`)' in js
    assert "cancelPointerSession('reflow')" in js
    assert "mode!=='smart'" in js and "mode==='guided'" in js
    assert "Smart mode will recompose" in js
    assert all(token in js for token in ("'compact'","'standard'","'prominent'","'hero'"))


def test_r14_authoring_ui_removes_persistent_shortcuts_and_preset_name_form():
    html=(PRODUCT/'assets/integrated_editor.html').read_text()
    assert 'Keyboard shortcuts' not in html
    assert 'Preset name' not in html
    assert 'Save as preset…' in html
    assert 'data-library-tab="elements"' in html and 'data-library-tab="presets"' in html
    assert 'Search 248 elements' in html


def test_r14_report_chrome_is_explicit_and_destructive_action_is_named():
    page=(PRODUCT/'page.py').read_text()
    assert "label='Report title'" in page and "label='Reports'" in page
    assert "ui.button('New report'" in page
    assert "ui.button('Import…'" in page
    assert "ui.button('Move to trash'" in page and 'color=negative' in page
    assert "Clean up empty reports" in page


def test_r14_normal_authoring_ui_does_not_expose_golden_connector_release_jargon():
    js=(PRODUCT/'assets/integrated_editor.mjs').read_text()
    assert 'Golden Connector' not in js
    assert 'Connections are routed automatically and stay editable.' in js


def test_r14_golden_connector_is_byte_identical():
    target=PRODUCT/'vendor/production_core/core/GOLDEN_CONNECTOR_ENGINE_V5_FROZEN.js'
    assert hashlib.sha256(target.read_bytes()).hexdigest()==GOLDEN


def test_r14_visualizer_runtime_dependencies_are_declared_by_the_release_package():
    requirements=[line.strip() for line in (ROOT/'requirements.txt').read_text(encoding='utf-8').splitlines() if line.strip() and not line.lstrip().startswith('#')]
    assert requirements==['nicegui==3.15.0','Pillow==12.3.0','python-pptx==1.0.2']
    project=tomllib.loads((ROOT/'pyproject.toml').read_text(encoding='utf-8'))['project']
    assert project['dependencies']==requirements
    test_requirements=(ROOT/'requirements-test.txt').read_text(encoding='utf-8')
    assert '-r requirements.txt' in test_requirements
    assert 'python-pptx' not in test_requirements and 'Pillow' not in test_requirements


def test_r14_production_dependency_import_contract():
    for package, expected, module in [('nicegui','3.15.0','nicegui'),('Pillow','12.3.0','PIL'),('python-pptx','1.0.2','pptx')]:
        assert version(package)==expected
        assert importlib.import_module(module)
    for module in ('company_ui.products.visualizer.cli','company_ui.products.visualizer.files','company_ui.products.visualizer.ppt_service'):
        assert importlib.import_module(module)


def _descriptions(prs: Presentation) -> list[str]:
    out=[]
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                nodes=shape._element.xpath('.//p:cNvPr')
                if nodes: out.append(nodes[0].get('descr') or '')
            except Exception:
                pass
    return out


def test_r14_ppt_export_embeds_exact_semantic_payload_for_zero_null_and_string_zero():
    model={'items':[
        {'id':'m','element':'Hero KPI','engine':'MetricEngine','title':'Metric zero','value':0,'unit':'hr'},
        {'id':'c','element':'Line Chart','engine':'CoreChartEngine','title':'Chart','rows':[{'label':'zero','value':0},{'label':'missing','value':None}]},
        {'id':'t','element':'Clean Table','engine':'TableEngine','title':'Table','customTable':{'headers':['kind','value'],'rows':[['numeric',0],['string','0'],['missing',None],['blank','']]}},
        {'id':'tm','element':'Event Timeline','engine':'TimelineEngine','title':'Timeline','milestones':[{'label':'Sequence only','date':None}]},
    ]}
    output=export_pptx(_pptx(),model); prs=Presentation(io.BytesIO(output)); descriptions='\n'.join(_descriptions(prs))
    for entry in model['items']:
        payload='VisualizerSemantic:'+json.dumps(entry,ensure_ascii=False,separators=(',',':'))
        assert payload in descriptions
    assert '"value":0' in descriptions and '["string","0"]' in descriptions and '["missing",null]' in descriptions
    assert '"date":null' in descriptions


def test_r14_ppt_export_uses_actual_metric_chart_and_table_values():
    model={'items':[
        {'id':'m','element':'Hero KPI','engine':'MetricEngine','title':'Metric zero','value':0,'unit':'hr'},
        {'id':'c','element':'Line Chart','engine':'CoreChartEngine','title':'Chart','rows':[{'label':'A','value':0},{'label':'B','value':5}]},
        {'id':'t','element':'Clean Table','engine':'TableEngine','title':'Table','customTable':{'headers':['Field','Value'],'rows':[['Zero',0],['Missing',None],['String zero','0']]}},
    ]}
    prs=Presentation(io.BytesIO(export_pptx(_pptx(),model))); slide=prs.slides[0]
    assert any(getattr(s,'has_text_frame',False) and 'Metric zero' in s.text and '\n0\n' in f'\n{s.text}\n' for s in slide.shapes)
    charts=[s.chart for s in slide.shapes if getattr(s,'has_chart',False)]
    assert len(charts)==1
    categories=[c.label for c in charts[0].plots[0].categories]
    assert categories==['A','B']
    tables=[s.table for s in slide.shapes if getattr(s,'has_table',False)]
    assert len(tables)==1
    values=[[cell.text for cell in row.cells] for row in tables[0].rows]
    assert values[1]==['Zero','0'] and values[2]==['Missing',''] and values[3]==['String zero','0']


def test_r14_sequence_only_timeline_does_not_invent_dates_in_ppt():
    model={'items':[{'id':'tm','element':'Event Timeline','engine':'TimelineEngine','title':'Timeline','milestones':[{'label':'First','date':None},{'label':'Second','date':None}]}]}
    prs=Presentation(io.BytesIO(export_pptx(_pptx(),model))); texts='\n'.join(s.text for s in prs.slides[0].shapes if getattr(s,'has_text_frame',False))
    assert 'First' in texts and 'Second' in texts
    assert '202' not in texts


def test_r14_ppt_export_paginates_dense_reports_without_truncation():
    items=[{'id':f'm{i}','element':'Hero KPI','engine':'MetricEngine','title':f'M{i}','value':i} for i in range(13)]
    prs=Presentation(io.BytesIO(export_pptx(_pptx(),{'items':items})))
    assert len(prs.slides)==2
    assert sum(1 for description in _descriptions(prs) if description.startswith('VisualizerSemantic:'))==13
